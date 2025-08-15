#!/usr/bin/env python3
import sys, re, os
from pptx import Presentation
from markdownify import markdownify as md

pptx_path = sys.argv[1]
prs = Presentation(pptx_path)

base_name_raw = os.path.splitext(os.path.basename(pptx_path))[0]
base_safe = re.sub(r'\s+', '-', base_name_raw)
os.makedirs(base_safe, exist_ok=True)

with open(f"{base_safe}.md", "w", encoding="utf-8") as f:
    for idx, slide in enumerate(prs.slides, 1):
        f.write(f"## Slide {idx}\n\n")

        # ---------- 1. 普通文本 ----------
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                f.write(md(shape.text_frame.text.strip()) + "\n\n")

            # ---------- 2. 表格 ----------
            if shape.has_table:
                table = shape.table
                rows = list(table.rows)
                headers = [
                    md(cell.text.replace('\r', '').replace('\n', '/').strip())
                    for cell in rows[0].cells
                ]
                f.write("| " + " | ".join(headers) + " |\n")
                f.write("|" + "|".join(["---"] * len(rows[0].cells)) + "|\n")
                for row in rows[1:]:
                    cells = [
                        md(cell.text.replace('\r', '').replace('\n', '/').strip())
                        for cell in row.cells
                    ]
                    f.write("| " + " | ".join(cells) + " |\n")
                f.write("\n")

            # ---------- 3. 图片 ----------
            if shape.shape_type == 13:  # 13 = Picture
                img = shape.image
                ext = img.ext
                img_name = f"{base_safe}/slide_{idx}_{shape.shape_id}.{ext}"
                with open(img_name, "wb") as img_file:
                    img_file.write(img.blob)
                f.write(f"![image]({img_name})\n\n")

        # ---------- 4. 演讲者备注 ----------
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                f.write("**Speaker Notes:**\n\n")
                f.write(md(notes) + "\n\n")

        # ---------- 5. 批注（注释） ----------
        # python-pptx ≥0.6.21 支持 slide.comments
        if hasattr(slide, "comments") and slide.comments:
            f.write("**Comments:**\n\n")
            for cmt in slide.comments:
                author = cmt.author or "Unknown"
                text   = (cmt.text or "").replace('\r', '').replace('\n', '/')
                f.write(f"- **{author}**: {text}\n")
            f.write("\n")

print("Done →", f"{base_safe}.md")

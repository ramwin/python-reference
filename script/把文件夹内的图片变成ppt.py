#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Xiang Wang <ramwin@qq.com>


from pathlib import Path
import sys

import click
from pptx import Presentation


@click.command()
@click.option("--source", default="./")
@click.option("--target", default="all.pptx")
def main(source, target):
    files = list(Path(source).iterdir())

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for img_path in files:
        if img_path.suffix.lower() in [".jpg", ".jpeg", ".png"]:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(img_path), 0, 0)
    prs.save(target)


if __name__ == "__main__":
    main()
